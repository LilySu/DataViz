import sys
from PyQt5.QtWidgets import QApplication, QWidget, QPushButton, QInputDialog, QLineEdit, QFileDialog
from PyQt5.QtGui import QIcon
from PyQt5.QtCore import pyqtSlot
# from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
import argparse
from datetime import date


class App(QWidget):

    def __init__(self):
        super().__init__()
        self.title = 'PyQt5 button - pythonspot.com'
        self.left = 10
        self.top = 10
        self.width = 320
        self.height = 200
        self.initUI()

    def initUI(self):
        self.setWindowTitle(self.title)
        self.setGeometry(self.left, self.top, self.width, self.height)

        button = QPushButton('PyQt5 button', self)
        button.setToolTip('This s an example button')
        button.move(100, 70)
        button.clicked.connect(self.on_click)

        self.show()

    @pyqtSlot()
    def on_click(self):
        print('PyQt5 button click')

        class App(QWidget):

            def __init__(self):
                super().__init__()
                self.title = 'PyQt5 file dialogs - pythonspot.com'
                self.left = 10
                self.top = 10
                self.width = 640
                self.height = 480
                self.initUI()

            def initUI(self):
                self.setWindowTitle(self.title)
                self.setGeometry(self.left, self.top, self.width, self.height)

                self.openFileNameDialog()
                self.openFileNamesDialog()
                self.saveFileDialog()

                self.show()

            def openFileNameDialog(self):
                options = QFileDialog.Options()
                options |= QFileDialog.DontUseNativeDialog
                fileName, _ = QFileDialog.getOpenFileName(self, "QFileDialog.getOpenFileName()", "",
                                                          "All Files (*);;Python Files (*.py)", options=options)
                if fileName:
                    print(fileName)

            def openFileNamesDialog(self):
                options = QFileDialog.Options()
                options |= QFileDialog.DontUseNativeDialog
                files, _ = QFileDialog.getOpenFileNames(self, "QFileDialog.getOpenFileNames()", "",
                                                        "All Files (*);;Python Files (*.py)", options=options)
                if files:
                    print(files)

            def saveFileDialog(self):
                options = QFileDialog.Options()
                options |= QFileDialog.DontUseNativeDialog
                fileName, _ = QFileDialog.getSaveFileName(self, "QFileDialog.getSaveFileName()", "",
                                                          "All Files (*);;Text Files (*.txt)", options=options)
                if fileName:
                    print(fileName)

                    def parse_args():
                        """ Setup the input and output arguments for the script
                        Return the parsed input and output files
                        """
                        parser = argparse.ArgumentParser(description='Analyze powerpoint file structure')
                        parser.add_argument('infile',
                                            type=argparse.FileType('r'),
                                            help='Powerpoint file to be analyzed')
                        parser.add_argument('outfile',
                                            type=argparse.FileType('w'),
                                            help='Output powerpoint')
                        return parser.parse_args()

                    def create_ppt(input, output):
                        """ Take the input powerpoint file and use it as the template for the output
                        file.
                        """
                        prs = Presentation(input)
                        # Use the output from analyze_ppt to understand which layouts and placeholders
                        # to use
                        # Create a title slide first
                        title_slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        title.text = "Quarterly Report"
                        subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())
                        # Create the summary graph
                        graph_slide_layout = prs.slide_layouts[8]
                        slide = prs.slides.add_slide(graph_slide_layout)
                        title = slide.shapes.title
                        subtitle.text = 'Results consistent with last quarter'
                        slide = prs.slides.add_slide(prs.slide_layouts[2])
                        title = slide.shapes.title
                        top = Inches(1.5)
                        left = Inches(0.25)
                        width = Inches(9.25)
                        height = Inches(5.0)
                        prs.save(output)

                    if __name__ == "__main__":
                        args = parse_args()
                        create_ppt(args.infile.name, args.outfile.name)

        if __name__ == '__main__':
            app = QApplication(sys.argv)
            ex = App()
            sys.exit(app.exec_())



if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = App()
    sys.exit(app.exec_())
