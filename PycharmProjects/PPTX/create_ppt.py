"""
See http://pbpython.com/creating-powerpoint.html for details on this script
Requires https://python-pptx.readthedocs.org/en/latest/index.html

Example program showing how to read in Excel, process with pandas and
output to a PowerPoint file.
"""

from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
import argparse
from datetime import date



def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Analyze powerpoint file structure')
    parser.add_argument('infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file to be analyzed')
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint')
    return parser.parse_args()

def create_ppt(input, output):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Quarterly Report"
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())
    # Create the summary graph
    graph_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    subtitle.text = 'Results consistent with last quarter'
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    top = Inches(1.5)
    left = Inches(0.25)
    width = Inches(9.25)
    height = Inches(5.0)
    prs.save(output)


if __name__ == "__main__":
    args = parse_args()
    create_ppt(args.infile.name, args.outfile.name)
